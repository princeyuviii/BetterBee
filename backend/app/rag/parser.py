"""
BetterBee — Document Parsers.

Extracts raw text content from uploaded files depending on the file format (PDF, DOCX, Markdown, Text, Excel, PowerPoint).
"""

import io
from abc import ABC, abstractmethod
from typing import Any
import pypdf
import docx
import openpyxl
import pptx
import structlog

logger = structlog.get_logger(__name__)


class DocumentParser(ABC):
    """Abstract interface defining document parsing."""

    @abstractmethod
    def parse(self, content_bytes: bytes) -> tuple[str, list[dict[str, Any]]]:
        """
        Parse raw file bytes and extract text and chunk metadata.
        
        Returns:
            Tuple of (raw_extracted_text, list_of_page_or_slide_metadata)
        """
        pass


class PDFParser(DocumentParser):
    """Parser for PDF files using pypdf."""

    def parse(self, content_bytes: bytes) -> tuple[str, list[dict[str, Any]]]:
        logger.debug("Parsing PDF file")
        pdf_file = io.BytesIO(content_bytes)
        reader = pypdf.PdfReader(pdf_file)
        
        extracted_text_parts = []
        page_metadata = []

        for page_idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            extracted_text_parts.append(text)
            page_metadata.append({
                "page_number": page_idx + 1,
            })

        return "\n\n".join(extracted_text_parts), page_metadata


class DocxParser(DocumentParser):
    """Parser for Microsoft Word files using python-docx."""

    def parse(self, content_bytes: bytes) -> tuple[str, list[dict[str, Any]]]:
        logger.debug("Parsing DOCX file")
        docx_file = io.BytesIO(content_bytes)
        doc = docx.Document(docx_file)
        
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)
        
        # Docx does not have clear page numbers, so we return a single metadata block
        return text, [{"page_number": 1}]


class ExcelParser(DocumentParser):
    """Parser for Microsoft Excel files using openpyxl."""

    def parse(self, content_bytes: bytes) -> tuple[str, list[dict[str, Any]]]:
        logger.debug("Parsing XLSX file")
        excel_file = io.BytesIO(content_bytes)
        wb = openpyxl.load_workbook(excel_file, read_only=True, data_only=True)
        
        sheet_texts = []
        metadata = []

        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                # Format row as tab-separated values
                row_str = "\t".join([str(val) if val is not None else "" for val in row])
                if row_str.replace("\t", "").strip():
                    rows.append(row_str)
            
            sheet_text = "\n".join(rows)
            sheet_texts.append(f"--- Sheet: {sheet_name} ---\n{sheet_text}")
            metadata.append({
                "sheet_name": sheet_name,
                "sheet_index": sheet_idx + 1,
            })

        return "\n\n".join(sheet_texts), metadata


class PptxParser(DocumentParser):
    """Parser for Microsoft PowerPoint files using python-pptx."""

    def parse(self, content_bytes: bytes) -> tuple[str, list[dict[str, Any]]]:
        logger.debug("Parsing PPTX file")
        pptx_file = io.BytesIO(content_bytes)
        presentation = pptx.Presentation(pptx_file)
        
        slide_texts = []
        metadata = []

        for slide_idx, slide in enumerate(presentation.slides):
            slide_text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text)
            
            slide_text = "\n".join(slide_text_parts)
            slide_texts.append(slide_text)
            metadata.append({
                "slide_number": slide_idx + 1,
            })

        return "\n\n".join(slide_texts), metadata


class TextParser(DocumentParser):
    """Parser for raw TXT and Markdown files."""

    def parse(self, content_bytes: bytes) -> tuple[str, list[dict[str, Any]]]:
        logger.debug("Parsing plain text/markdown file")
        try:
            text = content_bytes.decode("utf-8")
        except UnicodeDecodeError:
            # Fallback to latin-1
            text = content_bytes.decode("latin-1")
            
        return text, [{"page_number": 1}]


class ParserFactory:
    """Factory to get the appropriate parser based on file extension."""

    @staticmethod
    def get_parser(file_type: str) -> DocumentParser:
        clean_type = file_type.lower().strip(".")
        
        if clean_type == "pdf":
            return PDFParser()
        elif clean_type in ("docx", "doc"):
            return DocxParser()
        elif clean_type in ("xlsx", "xls", "csv"):
            return ExcelParser()
        elif clean_type in ("pptx", "ppt"):
            return PptxParser()
        elif clean_type in ("md", "txt", "markdown", "json", "xml", "html"):
            return TextParser()
        else:
            logger.warning("Unsupported file type, falling back to text parser", file_type=file_type)
            return TextParser()
